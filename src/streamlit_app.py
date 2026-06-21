import os
import io
import base64
import threading
import requests
import streamlit as st
from langsmith import traceable
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document as LCDocument
from rank_bm25 import BM25Okapi
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_bytes
import pdfplumber
import platform
from dotenv import load_dotenv
from concurrent.futures import ThreadPoolExecutor, as_completed

# =====================================================================
# Environment
# =====================================================================
load_dotenv()

api_key  = os.getenv("Api_key")
hf_token = os.getenv("Huggingface_api_key")
if hf_token:
    os.environ["HF_TOKEN"] = hf_token

if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    POPPLER_PATH = r"poppler-26.02.0\Library\bin"
else:
    POPPLER_PATH = None

# =====================================================================
# Page Config
# =====================================================================
st.set_page_config(
    page_title="RAG Document Q&A",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)

# =====================================================================
# CSS
# =====================================================================
st.markdown("""
<style>
.block-container { max-width: 1000px; padding-top: 1rem; padding-bottom: 2rem; }
.main-header { text-align: center; margin-bottom: 2rem; }
.main-header h1 { margin-bottom: 0.2rem; }
.main-header p { color: gray; }
.expanded-query { font-size: 0.78rem; color: #888; margin-top: -0.5rem; margin-bottom: 0.5rem; }
</style>
""", unsafe_allow_html=True)

# =====================================================================
# API Key Guard
# =====================================================================
if not api_key:
    st.error("⚠️ `Api_key` not found. Add your Groq API key as a HuggingFace Space secret named `Api_key`.")
    st.stop()

# =====================================================================
# Session State
# =====================================================================
_defaults = {
    "messages":           [],
    "vector_store":       None,
    "processed_file":     None,
    "image_status":       "idle",
    "image_status_msg":   "",
    "vlm_lock":           None,
    # BM25
    "bm25_index":         None,
    "bm25_chunks":        [],
    "bm25_metadatas":     [],
    # Summary cache
    "doc_summary":        None,    # None = not generated, str = cached
    "summary_file_id":    None,    # which file the summary is for
    "summary_raw_text":   "",      # first N chars of doc for summarisation
}
for k, v in _defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

if st.session_state.vlm_lock is None:
    st.session_state.vlm_lock = threading.Lock()

# =====================================================================
# Cached Resources
# =====================================================================
@st.cache_resource
def load_model(key):
    return ChatGroq(
        api_key=key,
        model="llama-3.1-8b-instant",
        temperature=0.7,
        max_tokens=1024,
        streaming=True,
    )

@st.cache_resource
def load_embeddings():
    return HuggingFaceEmbeddings(
        model_name="BAAI/bge-base-en-v1.5",
        model_kwargs={"device": "cpu"},
    )

model = load_model(api_key)

# =====================================================================
# Prompt  — citation-aware
# =====================================================================
rag_prompt = ChatPromptTemplate.from_template(
    """You are a helpful assistant. Answer the user's question using ONLY the provided context.

CITATION RULE: Whenever you use information from the context, cite it inline as:
  According to <filename> (Page <N>), ...
or at the end of the sentence as (<filename>, Page <N>).
If the answer is not in the context, say "I don't know based on the uploaded document."

Previous Conversation:
{history}

Context:
{context}

Question:
{question}

Answer (with citations):"""
)

# =====================================================================
# History Builder
# =====================================================================
def build_history(messages, n=10):
    recent = messages[-n:]
    lines  = []
    for m in recent:
        role = "User" if m["role"] == "user" else "Assistant"
        lines.append(f"{role}: {m['content']}")
    return "\n".join(lines) if lines else "None"

# =====================================================================
# Feature 3 — Query Expansion
# =====================================================================
def expand_query(original_query: str) -> str:
    """
    Ask the LLM to rewrite the user's query to be more self-contained
    and specific, as if referring to an uploaded document.
    Fast — uses a short non-streaming call.
    """
    try:
        expansion_prompt = (
            "You are a search query optimizer. "
            "Rewrite the following user question to be more specific and self-contained, "
            "as if the user is querying an uploaded document. "
            "Return ONLY the rewritten question — no explanation, no quotes.\n\n"
            f"Original question: {original_query}\n\n"
            "Rewritten question:"
        )
        resp = model.invoke(expansion_prompt)
        expanded = resp.content.strip()
        # Safety: if expansion looks broken or too long, fall back
        if not expanded or len(expanded) > 300:
            return original_query
        return expanded
    except Exception:
        return original_query

# =====================================================================
# Feature 4 — Hybrid Search (FAISS + BM25 with RRF)
# =====================================================================
def reciprocal_rank_fusion(ranked_lists: list[list], k: int = 60) -> list:
    """
    Merge multiple ranked lists using Reciprocal Rank Fusion.
    Each item is identified by its index in bm25_chunks / page_content.
    Returns deduplicated list of (score, item) sorted descending.
    """
    scores: dict = {}
    for ranked in ranked_lists:
        for rank, item in enumerate(ranked):
            key = item if isinstance(item, str) else item.page_content
            scores[key] = scores.get(key, 0.0) + 1.0 / (k + rank + 1)
    return sorted(scores.items(), key=lambda x: x[1], reverse=True)


def hybrid_retrieve(query: str, top_k: int = 6) -> list:
    """
    1. FAISS MMR retrieval  → top_k * 2 candidates
    2. BM25 retrieval       → top_k * 2 candidates
    3. RRF fusion           → top_k final chunks
    Returns list of LCDocument-like objects with .page_content and .metadata.
    """
    vs        = st.session_state.vector_store
    bm25      = st.session_state.bm25_index
    chunks    = st.session_state.bm25_chunks
    metadatas = st.session_state.bm25_metadatas

    if vs is None:
        return []

    # ── FAISS ──
    faiss_retriever = vs.as_retriever(
        search_type="mmr",
        search_kwargs={"k": top_k * 2, "fetch_k": top_k * 6},
    )
    faiss_results = faiss_retriever.invoke(query)   # list of LCDocument

    # ── BM25 ──
    bm25_results = []
    if bm25 and chunks:
        tokenized_q   = query.lower().split()
        bm25_scores   = bm25.get_scores(tokenized_q)
        top_indices   = sorted(range(len(bm25_scores)),
                               key=lambda i: bm25_scores[i],
                               reverse=True)[: top_k * 2]
        for idx in top_indices:
            doc = LCDocument(
                page_content=chunks[idx],
                metadata=metadatas[idx],
            )
            bm25_results.append(doc)

    # ── RRF Fusion ──
    faiss_keys = [d.page_content for d in faiss_results]
    bm25_keys  = [d.page_content for d in bm25_results]
    fused      = reciprocal_rank_fusion([faiss_keys, bm25_keys])

    # Build lookup: page_content → LCDocument (prefer faiss for metadata)
    lookup: dict[str, LCDocument] = {}
    for doc in bm25_results + faiss_results:   # faiss overwrites → preferred
        lookup[doc.page_content] = doc

    final = []
    for content, _score in fused[:top_k]:
        if content in lookup:
            final.append(lookup[content])

    return final


# =====================================================================
# Feature 2 — Document Summary (posted as chat message)
# =====================================================================
def generate_summary_streamed(raw_text: str):
    """
    Stream a structured summary into a Streamlit placeholder.
    Returns the full summary string.
    """
    sample = raw_text[:4000]   # first ~4k chars is enough for a good summary

    summary_prompt = (
        "You are a document analyst. Given the following document excerpt, produce a structured summary.\n\n"
        "Format your response EXACTLY like this (use these headers):\n\n"
        "## 📄 Document Summary\n"
        "<2-3 sentence overview of what this document is about>\n\n"
        "## 🏷️ Key Topics\n"
        "<bullet list of main topics covered>\n\n"
        "## 🔍 Important Entities\n"
        "<bullet list of key people, organisations, products, dates, or numbers mentioned>\n\n"
        f"Document excerpt:\n{sample}\n\n"
        "Summary:"
    )

    placeholder = st.empty()
    full_text   = ""
    for chunk in model.stream(summary_prompt):
        if chunk.content:
            full_text += chunk.content
            placeholder.markdown(full_text + "▌")
    placeholder.markdown(full_text)
    return full_text


# =====================================================================
# VLM — Groq vision via REST
# =====================================================================
def describe_image_with_vlm(image_bytes: bytes, source: str, page: int):
    try:
        b64  = base64.b64encode(image_bytes).decode("utf-8")
        mime = "image/jpeg" if image_bytes[:3] == b"\xff\xd8\xff" else "image/png"

        payload = {
            "model": "meta-llama/llama-4-scout-17b-16e-instruct",
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "image_url",
                     "image_url": {"url": f"data:{mime};base64,{b64}"}},
                    {"type": "text",
                     "text": (
                         "Describe in detail what you see in this document image: "
                         "any diagrams, charts, tables, screenshots, figures, or visual information. "
                         "Be specific and thorough."
                     )},
                ],
            }],
            "max_tokens": 512,
            "temperature": 0.2,
        }

        resp = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json=payload,
            timeout=60,
        )
        if resp.status_code == 200:
            return resp.json()["choices"][0]["message"]["content"].strip()
        return None
    except Exception:
        return None


# =====================================================================
# Image Extraction from PDF
# =====================================================================
def extract_images_from_pdf(pdf_bytes: bytes, filename: str) -> list:
    """
    Smart image extraction — only sends genuinely visual content to VLM.

    Rules:
    1. Embedded images: only kept if the image covers >= MIN_IMAGE_AREA_RATIO
       of the page area (filters out tiny icons, bullets, decorative dots).
    2. Full-page fallback render: only triggered when the page has
       < MIN_TEXT_CHARS of extractable text AND no qualifying embedded images.
       This catches diagram-only pages while skipping text pages that happen
       to have decorative borders/lines (the root cause of the 63-image bug).
    """
    MIN_IMAGE_AREA_RATIO = 0.08   # image must cover >= 8% of page to matter
    MIN_TEXT_CHARS       = 120    # pages with more text than this are skipped for full-page render

    results = []
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            page_renders = None  # lazy — only rendered if truly needed

            for page_no, page in enumerate(pdf.pages, start=1):
                page_area     = (page.width or 1) * (page.height or 1)
                page_text     = page.extract_text() or ""
                has_good_text = len(page_text.strip()) >= MIN_TEXT_CHARS

                embedded      = page.images or []
                kept_images   = []

                for img_obj in embedded:
                    x0 = img_obj.get("x0", 0)
                    y0 = img_obj.get("y0", 0)
                    x1 = img_obj.get("x1", page.width)
                    y1 = img_obj.get("y1", page.height)

                    if x1 <= x0 or y1 <= y0:
                        continue  # degenerate bbox

                    img_area  = (x1 - x0) * (y1 - y0)
                    img_ratio = img_area / page_area

                    if img_ratio >= MIN_IMAGE_AREA_RATIO:
                        kept_images.append(img_obj)

                if kept_images:
                    # Crop and render only the qualifying embedded images
                    for img_obj in kept_images:
                        try:
                            x0 = img_obj.get("x0", 0);  y0 = img_obj.get("y0", 0)
                            x1 = img_obj.get("x1", page.width); y1 = img_obj.get("y1", page.height)
                            buf = io.BytesIO()
                            page.crop((x0, y0, x1, y1)).to_image(resolution=150).save(buf, format="PNG")
                            results.append({
                                "page": page_no,
                                "image_bytes": buf.getvalue(),
                                "source": filename,
                            })
                        except Exception:
                            continue

                elif not has_good_text:
                    # No qualifying embedded images AND very little text →
                    # likely a diagram/infographic page → render full page
                    if page_renders is None:
                        page_renders = convert_from_bytes(
                            pdf_bytes, dpi=120, poppler_path=POPPLER_PATH
                        )
                    if page_no - 1 < len(page_renders):
                        buf = io.BytesIO()
                        page_renders[page_no - 1].save(buf, format="PNG")
                        results.append({
                            "page": page_no,
                            "image_bytes": buf.getvalue(),
                            "source": filename,
                        })
                # else: page has good text and no big images → skip VLM entirely

    except Exception:
        pass

    return results


# =====================================================================
# Background VLM Worker
# =====================================================================
def run_vlm_in_background(image_items: list, vector_store, lock):
    try:
        st.session_state.image_status     = "processing"
        st.session_state.image_status_msg = f"🔍 Analyzing {len(image_items)} image(s) in background…"

        embedding_model = load_embeddings()
        new_texts, new_metas = [], []

        def process_one(item):
            desc = describe_image_with_vlm(item["image_bytes"], item["source"], item["page"])
            if desc:
                return {
                    "text": f"[Image / Diagram on page {item['page']} of {item['source']}]\n{desc}",
                    "meta": {"source": item["source"], "page": item["page"], "type": "image_description"},
                }
            return None

        with ThreadPoolExecutor(max_workers=3) as executor:
            for future in as_completed({executor.submit(process_one, i): i for i in image_items}):
                r = future.result()
                if r:
                    new_texts.append(r["text"])
                    new_metas.append(r["meta"])

        if new_texts:
            with lock:
                new_vs = FAISS.from_texts(texts=new_texts, embedding=embedding_model, metadatas=new_metas)
                vector_store.merge_from(new_vs)
                # Also add image descriptions to BM25
                st.session_state.bm25_chunks.extend(new_texts)
                st.session_state.bm25_metadatas.extend(new_metas)
                tokenized = [t.lower().split() for t in st.session_state.bm25_chunks]
                st.session_state.bm25_index = BM25Okapi(tokenized)

            st.session_state.image_status     = "done"
            st.session_state.image_status_msg = f"✅ {len(new_texts)} image(s) analyzed and indexed."
        else:
            st.session_state.image_status     = "done"
            st.session_state.image_status_msg = "ℹ️ No describable images found."
    except Exception as e:
        st.session_state.image_status     = "error"
        st.session_state.image_status_msg = f"⚠️ Image analysis failed: {e}"


# =====================================================================
# Text Extraction Helpers
# =====================================================================
def extract_pdf(uploaded_file):
    docs      = []
    pdf_bytes = uploaded_file.read()
    reader    = PdfReader(io.BytesIO(pdf_bytes))
    total_pages = len(reader.pages)
    page_images = None

    for page_no, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if len(page_text.strip()) > 50:
            docs.append({"source": uploaded_file.name, "page": page_no + 1, "text": page_text})
        else:
            if page_images is None:
                page_images = convert_from_bytes(pdf_bytes, dpi=150, poppler_path=POPPLER_PATH)
            try:
                ocr_text = pytesseract.image_to_string(page_images[page_no], config="--psm 6")
                if ocr_text.strip():
                    docs.append({"source": uploaded_file.name, "page": page_no + 1, "text": ocr_text})
            except Exception as e:
                st.warning(f"OCR failed on page {page_no + 1}: {e}")

    return docs, total_pages, pdf_bytes


def extract_docx(uploaded_file):
    doc  = Document(uploaded_file)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"source": uploaded_file.name, "page": 1, "text": text}] if text.strip() else []


def extract_pptx(uploaded_file):
    prs  = Presentation(uploaded_file)
    docs = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        text = "\n".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text.strip():
            docs.append({"source": uploaded_file.name, "page": slide_no, "text": text})
    return docs


def extract_txt(uploaded_file):
    text = uploaded_file.read().decode("utf-8", errors="ignore")
    return [{"source": uploaded_file.name, "page": 1, "text": text}] if text.strip() else []


# =====================================================================
# Sidebar
# =====================================================================
with st.sidebar:
    st.title("📂 Documents")

    uploaded_files = st.file_uploader(
        "Upload Documents",
        type=["pdf", "docx", "pptx", "txt"],
        accept_multiple_files=True,
    )

    if uploaded_files:
        file_id = "-".join(f"{f.name}-{f.size}" for f in uploaded_files)

        if st.session_state.processed_file != file_id:
            with st.spinner("Processing documents…"):
                all_documents = []
                total_pages   = 0
                pdf_bytes_map = {}

                for uploaded_file in uploaded_files:
                    ext = uploaded_file.name.split(".")[-1].lower()
                    try:
                        if ext == "pdf":
                            docs, pages, raw_bytes = extract_pdf(uploaded_file)
                            all_documents.extend(docs)
                            total_pages += pages
                            pdf_bytes_map[uploaded_file.name] = raw_bytes
                        elif ext == "docx":
                            all_documents.extend(extract_docx(uploaded_file))
                        elif ext == "pptx":
                            all_documents.extend(extract_pptx(uploaded_file))
                        elif ext == "txt":
                            all_documents.extend(extract_txt(uploaded_file))
                    except Exception as e:
                        st.warning(f"Could not process {uploaded_file.name}: {e}")

                if not all_documents:
                    st.error("No text could be extracted.")
                    st.stop()

                splitter   = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                all_chunks = []
                metadatas  = []

                for doc in all_documents:
                    chunks = splitter.split_text(doc["text"])
                    all_chunks.extend(chunks)
                    metadatas.extend([
                        {"source": doc["source"], "page": doc["page"], "type": "text"}
                        for _ in chunks
                    ])

                if not all_chunks:
                    st.error("Chunking produced no results.")
                    st.stop()

                try:
                    embedding_model = load_embeddings()
                except Exception as e:
                    st.error(f"Embedding error: {e}")
                    st.stop()

                try:
                    vector_store = FAISS.from_texts(
                        texts=all_chunks, embedding=embedding_model, metadatas=metadatas
                    )
                    st.session_state.vector_store = vector_store
                except Exception as e:
                    st.error(f"FAISS error: {e}")
                    st.stop()

                # ── Build BM25 index ──
                tokenized_corpus = [c.lower().split() for c in all_chunks]
                st.session_state.bm25_index    = BM25Okapi(tokenized_corpus)
                st.session_state.bm25_chunks   = list(all_chunks)
                st.session_state.bm25_metadatas = list(metadatas)

                # ── Cache raw text for summary ──
                combined_text = "\n\n".join(d["text"] for d in all_documents)
                st.session_state.summary_raw_text = combined_text
                st.session_state.doc_summary      = None   # reset for new file
                st.session_state.summary_file_id  = file_id

                st.session_state.processed_file  = file_id
                st.session_state.messages        = []
                st.session_state.image_status    = "idle"
                st.session_state.image_status_msg = ""

            st.success("✅ Documents indexed — start chatting!")

            sources = sorted({m["source"] for m in metadatas})
            for src in sources:
                st.write(f"• {src}")

            col1, col2, col3 = st.columns(3)
            col1.metric("Files",  len(uploaded_files))
            col2.metric("Pages",  total_pages)
            col3.metric("Chunks", len(all_chunks))

            # ── Launch VLM background thread ──
            if pdf_bytes_map:
                all_image_items = []
                for fname, raw_bytes in pdf_bytes_map.items():
                    all_image_items.extend(extract_images_from_pdf(raw_bytes, fname))

                if all_image_items:
                    t = threading.Thread(
                        target=run_vlm_in_background,
                        args=(all_image_items, st.session_state.vector_store, st.session_state.vlm_lock),
                        daemon=True,
                    )
                    t.start()
                    st.session_state.image_status     = "processing"
                    st.session_state.image_status_msg = f"🔍 Analyzing {len(all_image_items)} image(s) in background…"

        else:
            st.info("Using cached vector store")

    # ── VLM status badge ──
    status = st.session_state.image_status
    msg    = st.session_state.image_status_msg
    if status == "processing":
        st.info(msg)
    elif status == "done":
        st.success(msg)
    elif status == "error":
        st.warning(msg)

    st.divider()

    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state.messages = []
        st.rerun()

# =====================================================================
# Header
# =====================================================================
st.markdown("""
<div class="main-header">
    <h1>📄 RAG Document Q&A</h1>
    <p>Upload a document and ask questions — images, diagrams, citations included</p>
</div>
""", unsafe_allow_html=True)

# =====================================================================
# Chat History Display
# =====================================================================
if not st.session_state.messages:
    st.info("📎 Upload a document from the sidebar and start asking questions.")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])
        # Show expanded query caption under user messages if present
        if message["role"] == "user" and message.get("expanded"):
            st.caption(f"🔍 Expanded query: _{message['expanded']}_")

# =====================================================================
# Chat Input Row  — 📋 Summary button  +  chat_input side by side
# =====================================================================
doc_ready = st.session_state.vector_store is not None

# Layout: narrow button col | wide input col
btn_col, _ = st.columns([1, 11])

summary_clicked = False
with btn_col:
    if doc_ready:
        summary_clicked = st.button("📋", help="Generate Document Summary", use_container_width=True)

question = st.chat_input("Ask anything about your document, including diagrams and images…")

# =====================================================================
# Handle Summary Button
# =====================================================================
if summary_clicked and doc_ready:
    # Use cached summary if available for the same file
    if (
        st.session_state.doc_summary is not None
        and st.session_state.summary_file_id == st.session_state.processed_file
    ):
        cached = st.session_state.doc_summary
        st.session_state.messages.append({"role": "assistant", "content": cached})
        st.rerun()
    else:
        # Generate fresh summary and stream it into chat
        st.session_state.messages.append({
            "role": "assistant",
            "content": "⏳ Generating document summary…"
        })
        with st.chat_message("assistant"):
            summary_text = generate_summary_streamed(st.session_state.summary_raw_text)

        # Update the last message with real content
        st.session_state.messages[-1]["content"] = summary_text
        st.session_state.doc_summary   = summary_text
        st.session_state.summary_file_id = st.session_state.processed_file
        st.rerun()

# =====================================================================
# Handle Chat Question
# =====================================================================
if question:
    if not doc_ready:
        st.warning("⚠️ Please upload a document first.")
    else:
        # ── Feature 3: Query Expansion ──
        with st.spinner("Expanding query…"):
            expanded_query = expand_query(question)

        is_expanded = expanded_query.lower().strip() != question.lower().strip()

        # Append user message (store expanded query for caption display)
        user_msg = {"role": "user", "content": question}
        if is_expanded:
            user_msg["expanded"] = expanded_query
        st.session_state.messages.append(user_msg)

        with st.chat_message("user"):
            st.markdown(question)
            if is_expanded:
                st.caption(f"🔍 Expanded query: _{expanded_query}_")

        # ── Feature 4: Hybrid Retrieval ──
        with st.chat_message("assistant"):
            with st.spinner("Retrieving…"):
                relevant_chunks = hybrid_retrieve(
                    expanded_query if is_expanded else question
                )

                # Build citation-rich context  — include source + page in each block
                context_parts = []
                for c in relevant_chunks:
                    src  = c.metadata.get("source", "Unknown")
                    page = c.metadata.get("page", "?")
                    context_parts.append(
                        f"[Source: {src}, Page {page}]\n{c.page_content}"
                    )
                context      = "\n\n---\n\n".join(context_parts)
                history_text = build_history(st.session_state.messages[:-1])

                # ── Feature 5: Citation-aware prompt ──
                formatted_prompt = rag_prompt.format(
                    history=history_text,
                    context=context,
                    question=expanded_query if is_expanded else question,
                )

            # Stream answer
            placeholder = st.empty()
            answer = ""
            for chunk in model.stream(formatted_prompt):
                if chunk.content:
                    answer += chunk.content
                    placeholder.markdown(answer + "▌")
            placeholder.markdown(answer)

            # Sources expander
            with st.expander("📚 Sources Used"):
                for i, chunk in enumerate(relevant_chunks, start=1):
                    source = chunk.metadata.get("source", "Unknown")
                    page   = chunk.metadata.get("page", "?")
                    ctype  = chunk.metadata.get("type", "text")
                    icon   = "🖼️" if ctype == "image_description" else "📄"
                    label  = "Image / Diagram" if ctype == "image_description" else "Text"
                    st.markdown(f"**Chunk {i}** | {icon} `{source}` — Page {page} _{label}_")
                    st.write(chunk.page_content[:500])
                    st.divider()

        st.session_state.messages.append({"role": "assistant", "content": answer})
        st.rerun()